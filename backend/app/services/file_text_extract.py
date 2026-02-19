from __future__ import annotations

import re
import zipfile
from io import BytesIO
from pathlib import Path

import openpyxl
import xlrd
from pptx import Presentation
from pypdf import PdfReader


def decode_text_with_fallback(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def _strip_xml_tags(text: str) -> str:
    cleaned = re.sub(r"<[^>]+>", " ", text)
    return re.sub(r"\s+", " ", cleaned).strip()


def extract_text_from_file(file_name: str, file_bytes: bytes) -> str:
    suffix = Path(file_name).suffix.lower()
    if suffix in {".txt", ".md", ".markdown", ".json", ".csv", ".log", ".xml", ".yaml", ".yml"}:
        return decode_text_with_fallback(file_bytes)

    if suffix == ".pdf":
        reader = PdfReader(BytesIO(file_bytes))
        chunks: list[str] = []
        for page in reader.pages:
            chunks.append(page.extract_text() or "")
        return "\n".join(chunks).strip()

    if suffix == ".docx":
        with zipfile.ZipFile(BytesIO(file_bytes)) as zf:
            xml_payload = zf.read("word/document.xml").decode("utf-8", errors="ignore")
            return _strip_xml_tags(xml_payload)

    if suffix == ".pptx":
        presentation = Presentation(BytesIO(file_bytes))
        chunks: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                text = shape.text.strip()
                if text:
                    chunks.append(text)
        return "\n".join(chunks).strip()

    if suffix == ".xlsx":
        workbook = openpyxl.load_workbook(BytesIO(file_bytes), data_only=True)
        rows: list[str] = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell) for cell in row if cell is not None]
                if cells:
                    rows.append("\t".join(cells))
        return "\n".join(rows).strip()

    if suffix == ".xls":
        workbook = xlrd.open_workbook(file_contents=file_bytes)
        rows: list[str] = []
        for sheet in workbook.sheets():
            for row_idx in range(sheet.nrows):
                cells = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                cells = [cell for cell in cells if cell and cell != "None"]
                if cells:
                    rows.append("\t".join(cells))
        return "\n".join(rows).strip()

    return decode_text_with_fallback(file_bytes)


def safe_storage_name(name: str, fallback: str = "file") -> str:
    cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1F]', "_", name).strip().strip(".")
    if not cleaned:
        cleaned = fallback
    return cleaned[:180]
